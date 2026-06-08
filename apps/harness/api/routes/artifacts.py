import base64
from datetime import datetime
from typing import List, Optional
from fastapi import APIRouter, Depends, HTTPException, UploadFile, File
from fastapi.responses import Response
from pydantic import BaseModel
from sqlalchemy import select, delete
from sqlalchemy.ext.asyncio import AsyncSession

from core.auth import require_auth
from db.session import get_db
from db.models import Artifact
from ingest.pipeline import ingest_file

router = APIRouter()

# MIME type map for generated binary artifacts
_BINARY_MIME = {
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pdf":  "application/pdf",
}


# ─── Response models ──────────────────────────────────────────────────────────

class ArtifactVersionOut(BaseModel):
    version: int
    created_at: datetime
    size_bytes: int

    class Config:
        from_attributes = True


class ArtifactOut(BaseModel):
    id: str
    filename: str
    type: str
    source: str
    source_id: Optional[str]
    project_ref: Optional[str]
    tags: list
    size_bytes: int
    version: int
    parent_id: Optional[str]
    created_at: datetime

    class Config:
        from_attributes = True


class ArtifactDetailOut(ArtifactOut):
    content: Optional[str] = None


class CreateArtifactRequest(BaseModel):
    filename: str
    type: str                          # document | code | report | spreadsheet | transcript
    source: str                        # chat | agent_job | cron | meeting | upload
    source_id: Optional[str] = None
    content: str
    project_ref: Optional[str] = None
    tags: List[str] = []
    parent_id: Optional[str] = None   # set to re-version an existing artifact


# ─── Helpers ──────────────────────────────────────────────────────────────────

def _classify_type(filename: str) -> str:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext in {"py", "ts", "tsx", "js", "jsx", "go", "rs", "java", "rb", "sh", "diff", "patch"}:
        return "code"
    if ext in {"csv", "xlsx", "xls"}:
        return "spreadsheet"
    if ext in {"txt", "vtt", "srt"}:
        return "transcript"
    return "document"


# ─── Routes ───────────────────────────────────────────────────────────────────

@router.get("", response_model=List[ArtifactOut])
async def list_artifacts(
    type: Optional[str] = None,
    source: Optional[str] = None,
    user_id: str = Depends(require_auth),
    db: AsyncSession = Depends(get_db),
):
    q = select(Artifact).where(Artifact.user_id == user_id)
    if type:
        q = q.where(Artifact.type == type)
    if source:
        q = q.where(Artifact.source == source)
    q = q.order_by(Artifact.created_at.desc()).limit(200)
    result = await db.execute(q)
    return result.scalars().all()


@router.get("/{artifact_id}", response_model=ArtifactDetailOut)
async def get_artifact(
    artifact_id: str,
    user_id: str = Depends(require_auth),
    db: AsyncSession = Depends(get_db),
):
    result = await db.execute(
        select(Artifact).where(Artifact.id == artifact_id, Artifact.user_id == user_id)
    )
    artifact = result.scalar_one_or_none()
    if not artifact:
        raise HTTPException(status_code=404, detail="Artifact not found")
    return artifact


@router.get("/{artifact_id}/versions", response_model=List[ArtifactOut])
async def list_artifact_versions(
    artifact_id: str,
    user_id: str = Depends(require_auth),
    db: AsyncSession = Depends(get_db),
):
    """Return all versions of an artifact (the one with this id + any that share parent_id)."""
    # Find the root — could be this artifact or its parent
    result = await db.execute(
        select(Artifact).where(Artifact.id == artifact_id, Artifact.user_id == user_id)
    )
    artifact = result.scalar_one_or_none()
    if not artifact:
        raise HTTPException(status_code=404, detail="Artifact not found")

    root_id = artifact.parent_id or artifact_id

    versions_result = await db.execute(
        select(Artifact)
        .where(
            Artifact.user_id == user_id,
            (Artifact.id == root_id) | (Artifact.parent_id == root_id),
        )
        .order_by(Artifact.version.desc())
    )
    return versions_result.scalars().all()


@router.post("", response_model=ArtifactDetailOut, status_code=201)
async def create_artifact(
    body: CreateArtifactRequest,
    user_id: str = Depends(require_auth),
    db: AsyncSession = Depends(get_db),
):
    artifact_type = body.type or _classify_type(body.filename)

    # If this is a new version of an existing artifact, bump version number
    version = 1
    if body.parent_id:
        result = await db.execute(
            select(Artifact)
            .where(
                Artifact.user_id == user_id,
                (Artifact.id == body.parent_id) | (Artifact.parent_id == body.parent_id),
            )
            .order_by(Artifact.version.desc())
            .limit(1)
        )
        latest = result.scalar_one_or_none()
        if latest:
            version = latest.version + 1

    artifact = Artifact(
        user_id=user_id,
        filename=body.filename,
        type=artifact_type,
        source=body.source,
        source_id=body.source_id,
        content=body.content,
        version=version,
        parent_id=body.parent_id,
        project_ref=body.project_ref,
        tags=body.tags,
        size_bytes=len(body.content.encode("utf-8")),
    )
    db.add(artifact)
    await db.commit()
    await db.refresh(artifact)
    return artifact


@router.delete("/{artifact_id}", status_code=204)
async def delete_artifact(
    artifact_id: str,
    user_id: str = Depends(require_auth),
    db: AsyncSession = Depends(get_db),
):
    result = await db.execute(
        select(Artifact).where(Artifact.id == artifact_id, Artifact.user_id == user_id)
    )
    artifact = result.scalar_one_or_none()
    if not artifact:
        raise HTTPException(status_code=404, detail="Artifact not found")
    await db.execute(delete(Artifact).where(Artifact.id == artifact_id))
    await db.commit()


@router.get("/{artifact_id}/download")
async def download_artifact(
    artifact_id: str,
    user_id: str = Depends(require_auth),
    db: AsyncSession = Depends(get_db),
):
    """
    Serve an artifact as a downloadable binary file.
    Text artifacts are served as text/plain; base64-encoded binary artifacts
    (DOCX, PPTX, PDF) are decoded and served with the correct MIME type.
    """
    result = await db.execute(
        select(Artifact).where(Artifact.id == artifact_id, Artifact.user_id == user_id)
    )
    art = result.scalar_one_or_none()
    if not art:
        raise HTTPException(status_code=404, detail="Artifact not found")

    filename = art.filename or "download"
    ext = ("." + filename.rsplit(".", 1)[-1].lower()) if "." in filename else ""

    if art.content and art.content.startswith("base64:"):
        # Binary artifact — decode and serve
        raw = base64.b64decode(art.content[7:])
        mime = _BINARY_MIME.get(ext, "application/octet-stream")
        return Response(
            content=raw,
            media_type=mime,
            headers={"Content-Disposition": f'attachment; filename="{filename}"'},
        )

    # Text artifact — serve as UTF-8
    content = art.content or ""
    return Response(
        content=content.encode("utf-8"),
        media_type="text/plain; charset=utf-8",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.get("/{artifact_id}/view")
async def view_artifact(
    artifact_id: str,
    user_id: str = Depends(require_auth),
    db: AsyncSession = Depends(get_db),
):
    """
    Serve artifact inline (for iframe/preview embedding).
    Same as /download but uses Content-Disposition: inline so the browser
    renders it in-place rather than triggering a save dialog.
    """
    result = await db.execute(
        select(Artifact).where(Artifact.id == artifact_id, Artifact.user_id == user_id)
    )
    art = result.scalar_one_or_none()
    if not art:
        raise HTTPException(status_code=404, detail="Artifact not found")

    filename = art.filename or "preview"
    ext = ("." + filename.rsplit(".", 1)[-1].lower()) if "." in filename else ""

    if art.content and art.content.startswith("base64:"):
        raw = base64.b64decode(art.content[7:])
        mime = _BINARY_MIME.get(ext, "application/octet-stream")
        return Response(
            content=raw,
            media_type=mime,
            headers={"Content-Disposition": f'inline; filename="{filename}"'},
        )

    content = art.content or ""
    return Response(
        content=content.encode("utf-8"),
        media_type="text/plain; charset=utf-8",
        headers={"Content-Disposition": f'inline; filename="{filename}"'},
    )


@router.get("/{artifact_id}/preview")
async def preview_artifact(
    artifact_id: str,
    user_id: str = Depends(require_auth),
    db: AsyncSession = Depends(get_db),
):
    """
    Extract human-readable preview text from binary artifacts.
    Returns JSON { text: str, type: "docx"|"pptx"|"text" }.
    PDF files should use /view (rendered by the browser via iframe).
    """
    result = await db.execute(
        select(Artifact).where(Artifact.id == artifact_id, Artifact.user_id == user_id)
    )
    art = result.scalar_one_or_none()
    if not art:
        raise HTTPException(status_code=404, detail="Artifact not found")

    content = art.content or ""
    filename = (art.filename or "").lower()

    if not content.startswith("base64:"):
        return {"text": content[:8000], "type": "text"}

    raw = base64.b64decode(content[7:])

    if filename.endswith(".pdf"):
        try:
            from ingest.parsers import pdf as _pdf_parser
            text = _pdf_parser.extract(raw)
            return {"text": text[:8000], "type": "text"}
        except Exception as exc:
            return {"text": f"(preview extraction failed: {exc})", "type": "text"}

    if filename.endswith(".docx"):
        try:
            import docx as _docx
            from io import BytesIO as _BIO
            doc = _docx.Document(_BIO(raw))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return {"text": "\n\n".join(paragraphs[:300]), "type": "docx"}
        except Exception as exc:
            return {"text": f"(preview extraction failed: {exc})", "type": "text"}

    if filename.endswith(".pptx"):
        try:
            from pptx import Presentation as _Prs
            from io import BytesIO as _BIO
            prs = _Prs(_BIO(raw))
            slides_text = []
            for i, slide in enumerate(prs.slides):
                texts = [
                    shp.text_frame.text.strip()
                    for shp in slide.shapes
                    if hasattr(shp, "text_frame") and shp.text_frame.text.strip()
                ]
                if texts:
                    slides_text.append(f"## Slide {i + 1}\n" + "\n".join(texts))
            return {"text": "\n\n".join(slides_text), "type": "pptx"}
        except Exception as exc:
            return {"text": f"(preview extraction failed: {exc})", "type": "text"}

    # Unknown binary (shouldn't normally reach here for DOCX/PPTX/PDF)
    return {"text": "", "type": "binary"}


@router.post("/ingest", response_model=ArtifactDetailOut, status_code=201)
async def ingest_artifact(
    file: UploadFile = File(...),
    user_id: str = Depends(require_auth),
    db: AsyncSession = Depends(get_db),
):
    """
    Upload any file and have it parsed into an Artifact.
    Binary office/PDF files are stored as base64 so the original can be
    downloaded and previewed. Text/code files are stored as extracted text.
    """
    content_bytes = await file.read()
    filename = file.filename or "upload"
    mime_type = file.content_type or ""

    if not content_bytes:
        raise HTTPException(status_code=400, detail="Empty file")

    from ingest.pipeline import detect_mime, _artifact_type_from_ext

    detected_mime = mime_type or detect_mime(filename, content_bytes)
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    artifact_type = _artifact_type_from_ext(filename)

    # HEIC/HEIF photos: reject at upload time with a clear message
    if ext in ("heic", "heif") or detected_mime in ("image/heic", "image/heif"):
        raise HTTPException(
            status_code=422,
            detail=(
                "HEIC photos cannot be analyzed directly. "
                "Please export the photo as JPEG or PNG (in your Photos app: Share → Save as JPEG) and try again."
            ),
        )

    # BMP/TIFF: image-like but Anthropic Vision doesn't accept them
    if ext in ("bmp", "tiff", "tif") or detected_mime in ("image/bmp", "image/tiff", "image/x-tiff", "image/x-bmp"):
        raise HTTPException(
            status_code=422,
            detail=(
                f"'{filename}' is in {ext.upper()} format which cannot be analyzed. "
                "Please convert to JPEG or PNG and try again."
            ),
        )

    # Binary formats are stored as base64 — preserves the original file for
    # download/preview. Text extraction happens on-the-fly when injected into chat.
    # Images are also stored as base64 so Sonnet can analyze them directly via vision.
    _BINARY_EXTS = {"pdf", "docx", "ppt", "pptx", "xls", "xlsx"}
    _BINARY_MIMES = {
        "application/pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.ms-powerpoint",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-excel",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    }
    _IMAGE_EXTS = {"jpg", "jpeg", "png", "gif", "webp"}
    _IMAGE_MIMES = {"image/jpeg", "image/jpg", "image/png", "image/gif", "image/webp"}
    # Legacy .doc: stored as base64 (no text extraction), with a note when injected into chat
    _DOC_LEGACY_EXTS = {"doc"}
    _DOC_LEGACY_MIMES = {"application/msword"}

    if (ext in _BINARY_EXTS or detected_mime in _BINARY_MIMES
            or ext in _IMAGE_EXTS or detected_mime in _IMAGE_MIMES
            or ext in _DOC_LEGACY_EXTS or detected_mime in _DOC_LEGACY_MIMES):
        content = "base64:" + base64.b64encode(content_bytes).decode()
        size_bytes = len(content_bytes)
    else:
        result = await ingest_file(
            content_bytes=content_bytes,
            filename=filename,
            mime_type=detected_mime,
        )
        content = result["content"]
        artifact_type = result["artifact_type"]
        size_bytes = len(content.encode("utf-8"))

    artifact = Artifact(
        user_id=user_id,
        filename=filename,
        type=artifact_type,
        source="upload",
        content=content,
        version=1,
        size_bytes=size_bytes,
        tags=[],
    )
    db.add(artifact)
    await db.commit()
    await db.refresh(artifact)
    return artifact
