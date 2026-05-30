import asyncio
import base64
import logging
from datetime import datetime, timezone
from typing import Any, Dict, List, Optional
import json as _json
from fastapi import APIRouter, Depends, File, Form, HTTPException, UploadFile
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from sqlalchemy import select, desc, update as sa_update, delete as sa_delete, func
from sqlalchemy.ext.asyncio import AsyncSession

log = logging.getLogger(__name__)

from core.auth import require_auth
from core.router import classify
from core.model_client import (
    get_model_client, ModelClient, ModelTier,
    PROPOSE_CALENDAR_EVENT_TOOL, PROPOSE_TASK_TOOL,
    CREATE_TASK_TOOL, CREATE_CALENDAR_EVENT_TOOL,
    SAVE_MEMORY_TOOL, SAVE_TO_SECOND_BRAIN_TOOL,
    READ_EMAIL_TOOL, SEND_EMAIL_TOOL, READ_MEETING_TOOL, SYNC_MEETINGS_TOOL, WEB_SEARCH_TOOL,
    GENERATE_DOCUMENT_TOOL, GENERATE_PRESENTATION_TOOL, GENERATE_PDF_TOOL,
    LOOKUP_CONTACT_TOOL, SEARCH_CONTACTS_TOOL,
    CREATE_CONTACT_TOOL, UPDATE_CONTACT_TOOL,
    SEARCH_PLACES_TOOL, SAVE_PLACE_TOOL, GET_SAVED_PLACES_TOOL,
    CREATE_AGENT_JOB_TOOL,
)
from core.context_assembler import assemble
from core.streaming import sse_event, sse_done
from db.session import get_db, AsyncSessionLocal
from db.models import Conversation, Message, User, Task, Artifact

router = APIRouter()

# Strong references so background tasks aren't GC'd before they finish
_active_bg_tasks: set = set()


# ── Schemas ───────────────────────────────────────────────────────────────────

class ConversationOut(BaseModel):
    id: str
    title: Optional[str]
    created_at: datetime

    class Config:
        from_attributes = True

    def model_post_init(self, __context) -> None:
        # Cap title at 60 chars regardless of what's stored in DB
        if self.title and len(self.title) > 60:
            object.__setattr__(self, "title", self.title[:60])


class MessageOut(BaseModel):
    id: str
    role: str
    content: str
    model_used: Optional[str]
    tokens_used: int
    tool_results: List[dict] = []
    created_at: datetime

    class Config:
        from_attributes = True


class ConversationDetailOut(BaseModel):
    id: str
    title: Optional[str]
    created_at: datetime
    messages: List[MessageOut]

    class Config:
        from_attributes = True

    def model_post_init(self, __context) -> None:
        if self.title and len(self.title) > 60:
            object.__setattr__(self, "title", self.title[:60])


_IMAGE_TYPES = {"image/jpeg", "image/png", "image/gif", "image/webp"}


def _xlsx_to_text(data: bytes) -> str:
    from io import BytesIO
    import openpyxl
    wb = openpyxl.load_workbook(BytesIO(data), read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue
        parts.append(f"Sheet: {sheet.title}")
        for row in rows:
            parts.append("\t".join("" if v is None else str(v) for v in row))
    return "\n".join(parts)


def _xls_to_text(data: bytes) -> str:
    import xlrd
    wb = xlrd.open_workbook(file_contents=data)
    parts = []
    for sheet in wb.sheets():
        parts.append(f"Sheet: {sheet.name}")
        for r in range(sheet.nrows):
            parts.append("\t".join(str(sheet.cell_value(r, c)) for c in range(sheet.ncols)))
    return "\n".join(parts)


def _pptx_to_text(data: bytes) -> str:
    from io import BytesIO
    from pptx import Presentation
    prs = Presentation(BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [s.text for s in slide.shapes if s.has_text_frame and s.text.strip()]
        if texts:
            parts.append(f"Slide {i}: " + " | ".join(texts))
    return "\n".join(parts)


async def _process_attachment(upload: UploadFile) -> dict:
    data = await upload.read()
    ct = (upload.content_type or "").lower()
    fname = upload.filename or "file"

    if ct.startswith("image/"):
        media_type = ct if ct in _IMAGE_TYPES else "image/jpeg"
        return {
            "kind": "image",
            "block": {
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": media_type,
                    "data": base64.standard_b64encode(data).decode(),
                },
            },
        }

    text = ""
    try:
        if ct == "application/pdf" or fname.endswith(".pdf"):
            import fitz
            doc = fitz.open(stream=data, filetype="pdf")
            text = "\n".join(page.get_text() for page in doc)

        elif "wordprocessingml" in ct or fname.endswith(".docx"):
            from io import BytesIO
            import docx as _docx
            document = _docx.Document(BytesIO(data))
            text = "\n".join(p.text for p in document.paragraphs if p.text.strip())

        elif "spreadsheetml" in ct or fname.endswith(".xlsx"):
            text = _xlsx_to_text(data)

        elif ct == "application/vnd.ms-excel" or fname.endswith(".xls"):
            text = _xls_to_text(data)

        elif "presentationml" in ct or fname.endswith(".pptx"):
            text = _pptx_to_text(data)

        elif ct in ("text/csv", "application/csv") or fname.endswith(".csv"):
            text = data.decode("utf-8", errors="replace")

        elif ct.startswith("text/") or fname.endswith((".txt", ".md", ".json", ".yaml", ".yml", ".xml")):
            text = data.decode("utf-8", errors="replace")

        else:
            decoded = data.decode("utf-8", errors="replace")
            replacement_ratio = decoded.count("�") / max(len(decoded), 1)
            if replacement_ratio > 0.1:
                return {"kind": "unsupported", "filename": fname, "mime": ct}
            text = decoded

    except Exception as exc:
        log.warning("Attachment extraction failed (%s): %s", fname, exc)
        text = f"[Extraction failed: {exc}]"

    return {"kind": "text", "filename": fname, "text": text}


# ── Helpers ───────────────────────────────────────────────────────────────────

async def _get_or_create_user(user_id: str, db: AsyncSession) -> User:
    result = await db.execute(select(User).where(User.id == user_id))
    user = result.scalar_one_or_none()
    if not user:
        user = User(id=user_id, name="Mike Villar")
        db.add(user)
        await db.flush()
    return user


async def _extract_and_save_facts(
    db,
    user_id: str,
    user_content: str,
    assistant_content: str,
    client: ModelClient,
) -> None:
    """
    Extract personal facts from what the USER said and save as memories.
    Only uses user_content as the fact source — assistant responses are NOT facts about Mike.
    Each fact is domain-tagged and saved individually for clean retrieval.
    """
    if len(user_content) < 20:
        return

    try:
        resp = await client.anthropic.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=300,
            system=(
                "Extract personal facts about Mike from his message. Only use what HE said — "
                "ignore the assistant context completely.\n"
                "Output each fact on its own line in this exact format:\n"
                "  DOMAIN|fact in third person\n"
                "Valid domains: work, personal, health, cycling, client\n"
                "Examples:\n"
                "  cycling|Mike's BMC Roadmachine has Ultegra Di2 and weighs 7.2 kg\n"
                "  health|Mike has Maxicare insurance, premium due June 1\n"
                "  work|Mike decided to delay the OpenRice campaign until Q3\n"
                "  client|Mike's NCH Inc. contact is Jaime Santos\n"
                "Output SKIP if no personal facts about Mike are present."
            ),
            messages=[{"role": "user", "content": f"Mike's message: {user_content[:500]}"}],
        )
        raw = resp.content[0].text.strip()
        if not raw or raw.upper() == "SKIP":
            return

        from memory import mnemon as _mnemon
        valid_domains = {"work", "personal", "health", "cycling", "client"}

        for line in raw.splitlines():
            line = line.strip().lstrip("•-* ")
            if not line or line.upper() == "SKIP":
                continue
            if "|" in line:
                domain, _, fact = line.partition("|")
                domain = domain.strip().lower()
                fact = fact.strip()
                if domain not in valid_domains:
                    domain = "work"
            else:
                fact = line
                domain = "work"

            if fact and len(fact) > 10:
                await _mnemon.write(
                    db, user_id, fact,
                    domain=domain, source="conversation", importance=4,
                )
    except Exception as exc:
        log.warning("Memory fact extraction failed: %s", exc)



def _strip_tool_artifacts(text: str) -> str:
    """Remove raw tool-call tags and internal tool syntax from message text."""
    import re
    # Strip XML-style tool call blocks: <tool_name>...</tool_name> or <save_memory>...</save_memory>
    text = re.sub(r"<[a-z_]+>.*?</[a-z_]+>", "", text, flags=re.DOTALL)
    # Strip standalone opening/closing tags: <save_memory>, </save_memory>
    text = re.sub(r"</?[a-z_]+>", "", text)
    # Strip function-call style: save_memory(content='...') or tool_name(...)
    text = re.sub(r"\b[a-z_]+\([^)]{0,200}\)", "", text)
    return text.strip()


async def _generate_title(messages: list, client: ModelClient) -> Optional[str]:
    """Generate a 3-5 word conversation title from recent exchanges."""
    recent = messages[-8:]
    context = "\n".join(
        f"{m['role'].upper()}: {_strip_tool_artifacts(str(m.get('content', '')))[:300]}"
        for m in recent
        if m.get("content") and not str(m.get("content", "")).startswith("<tool")
    )
    try:
        resp = await client.anthropic.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=15,
            system="Write a 3-5 word title for this conversation. No quotes, no punctuation, no explanation. Just the title.",
            messages=[{"role": "user", "content": context}],
        )
        return resp.content[0].text.strip()[:60]
    except Exception as exc:
        log.warning("Title generation failed: %s", exc)
        return None


async def _get_conversation(conv_id: str, user_id: str, db: AsyncSession) -> Conversation:
    result = await db.execute(
        select(Conversation).where(
            Conversation.id == conv_id,
            Conversation.user_id == user_id,
        )
    )
    conv = result.scalar_one_or_none()
    if not conv:
        raise HTTPException(status_code=404, detail="Conversation not found")
    return conv


# ── Routes ────────────────────────────────────────────────────────────────────

@router.get("/conversations", response_model=List[ConversationOut])
async def list_conversations(
    user_id: str = Depends(require_auth),
    db: AsyncSession = Depends(get_db),
):
    # Sort by most recent activity: join to messages and order by max(message.created_at),
    # falling back to conversation.created_at for new conversations with no messages yet.
    latest_msg = (
        select(Message.conversation_id, func.max(Message.created_at).label("last_at"))
        .group_by(Message.conversation_id)
        .subquery()
    )
    result = await db.execute(
        select(Conversation)
        .outerjoin(latest_msg, Conversation.id == latest_msg.c.conversation_id)
        .where(Conversation.user_id == user_id)
        .order_by(desc(func.coalesce(latest_msg.c.last_at, Conversation.created_at)))
        .limit(50)
    )
    return result.scalars().all()


@router.post("/conversations", response_model=ConversationOut, status_code=201)
async def create_conversation(
    user_id: str = Depends(require_auth),
    db: AsyncSession = Depends(get_db),
):
    await _get_or_create_user(user_id, db)
    conv = Conversation(user_id=user_id)
    db.add(conv)
    await db.commit()
    await db.refresh(conv)
    return conv


@router.get("/conversations/{conversation_id}", response_model=ConversationDetailOut)
async def get_conversation(
    conversation_id: str,
    user_id: str = Depends(require_auth),
    db: AsyncSession = Depends(get_db),
):
    conv = await _get_conversation(conversation_id, user_id, db)
    result = await db.execute(
        select(Message)
        .where(Message.conversation_id == conversation_id)
        .order_by(Message.created_at)
    )
    messages = result.scalars().all()

    title = conv.title

    # Lazily truncate titles that predate the 60-char cap
    if title and len(title) > 60:
        title = title[:60]
        await db.execute(
            sa_update(Conversation)
            .where(Conversation.id == conversation_id)
            .values(title=title)
        )
        await db.commit()

    if not title and messages:
        msg_dicts = [{"role": m.role, "content": m.content} for m in messages]
        client = get_model_client()
        title = await _generate_title(msg_dicts, client)
        if title:
            await db.execute(
                sa_update(Conversation)
                .where(Conversation.id == conversation_id)
                .values(title=title)
            )
            await db.commit()

    return ConversationDetailOut(
        id=conv.id,
        title=title,
        created_at=conv.created_at,
        messages=messages,
    )


@router.delete("/conversations/{conversation_id}", status_code=204)
async def delete_conversation(
    conversation_id: str,
    user_id: str = Depends(require_auth),
    db: AsyncSession = Depends(get_db),
):
    conv = await _get_conversation(conversation_id, user_id, db)
    # Delete messages first — no DB-level cascade on this FK
    await db.execute(sa_delete(Message).where(Message.conversation_id == conversation_id))
    await db.delete(conv)
    await db.commit()


@router.post("/conversations/{conversation_id}/messages")
async def send_message(
    conversation_id: str,
    content: str = Form(default=""),
    files: List[UploadFile] = File(default=[]),
    artifact_id: Optional[str] = Form(default=None),
    tier_override: Optional[str] = Form(default=None),
    location_lat: Optional[float] = Form(default=None),
    location_lng: Optional[float] = Form(default=None),
    user_id: str = Depends(require_auth),
    db: AsyncSession = Depends(get_db),
):
    """Stream an assistant reply via SSE. Saves both messages to DB after completion."""
    conv = await _get_conversation(conversation_id, user_id, db)

    # Process attachments
    image_blocks: List[Any] = []
    doc_snippets: List[str] = []
    for upload in files:
        result = await _process_attachment(upload)
        if result["kind"] == "image":
            image_blocks.append(result["block"])
        elif result["kind"] == "unsupported":
            doc_snippets.append(
                f"[Attached: {result['filename']} — file type '{result['mime']}' could not be extracted as text. "
                f"Let the user know what types are supported (PDF, DOCX, XLSX, XLS, PPTX, CSV, TXT, images).]"
            )
        else:
            doc_snippets.append(f"[Attached: {result['filename']}]\n{result['text']}")

    # If artifact_id provided, inject artifact content as invisible model context
    if artifact_id:
        art_result = await db.execute(
            select(Artifact).where(Artifact.id == artifact_id, Artifact.user_id == user_id)
        )
        art_obj = art_result.scalar_one_or_none()
        if art_obj and art_obj.content:
            MAX_ART = 12000
            art_text = art_obj.content
            if len(art_text) > MAX_ART:
                art_text = art_text[:MAX_ART] + f"\n\n[… truncated, {len(art_obj.content):,} total chars]"
            doc_snippets.append(
                f"[UPLOADED FILE: {art_obj.filename}]\n{art_text}\n\n"
                f"Analyze this file and provide a clear summary of the key insights."
            )

    # Classify the request to pick the right model tier.
    # Images always need vision (Claude). Override wins if provided.
    if image_blocks:
        tier = ModelTier.TIER3
    elif artifact_id:
        tier = ModelTier.TIER3  # artifact analysis always uses the frontier model
    elif tier_override:
        tier = ModelTier(tier_override)
    elif content:
        tier = await classify(content)
    elif doc_snippets:
        tier = ModelTier.TIER2
    else:
        tier = ModelTier.TIER2

    # Build the text that goes to the model (doc text prepended)
    full_text = ("\n\n".join(doc_snippets) + "\n\n" + content).strip() if doc_snippets else content

    # Store plain text in DB (don't persist image bytes)
    db_content = content or " · ".join(
        (f"[image]" if b["type"] == "image" else "") for b in image_blocks
    ) or "[attachment]"
    user_msg = Message(
        conversation_id=conversation_id,
        role="user",
        content=db_content,
    )
    db.add(user_msg)
    await db.commit()

    db_result = await db.execute(
        select(Message)
        .where(Message.conversation_id == conversation_id)
        .order_by(Message.created_at)
        .limit(40)
    )
    history = db_result.scalars().all()
    messages: List[Dict[str, Any]] = [
        {"role": m.role, "content": m.content} for m in history if m.role != "system"
    ]

    # ── Normalise for Anthropic API ──────────────────────────────────────────
    # Race condition: background tasks save the assistant reply asynchronously.
    # A new user message can arrive *before* the previous assistant reply is
    # committed, so the DB row for the assistant ends up with a later
    # created_at than the new user row.  That produces a history like:
    #   [user_old, user_new, assistant_old]
    # …which the API rejects as "assistant prefill" (last msg = assistant).
    # The same ordering bug can also leave two consecutive user messages.
    # Fix: strip trailing assistant messages, then merge any consecutive
    # same-role pairs so the final list strictly alternates and ends with user.
    def _sanitize(msgs: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        # 1. Drop trailing assistant messages.
        trimmed = list(msgs)
        while trimmed and trimmed[-1]["role"] != "user":
            trimmed.pop()
        if not trimmed:
            return trimmed
        # 2. Collapse consecutive same-role messages (merge text; keep newer on rich).
        result: List[Dict[str, Any]] = [dict(trimmed[0])]
        for m in trimmed[1:]:
            prev = result[-1]
            if prev["role"] == m["role"]:
                p, c = prev["content"], m["content"]
                if isinstance(p, str) and isinstance(c, str):
                    result[-1] = {"role": prev["role"], "content": p + "\n\n" + c}
                else:
                    result[-1] = dict(m)   # rich content: keep the newer block
            else:
                result.append(dict(m))
        return result

    messages = _sanitize(messages)

    # Replace last user message with rich content if needed
    if messages and messages[-1]["role"] == "user":
        if image_blocks:
            parts: List[Any] = []
            if full_text:
                parts.append({"type": "text", "text": full_text})
            parts.extend(image_blocks)
            messages[-1]["content"] = parts
        elif full_text != content:
            messages[-1]["content"] = full_text

    client = get_model_client()
    # When RunPod (Tier2) is cold and falls back to Claude, upgrade to Tier3 so the
    # request gets full tools + capabilities block instead of a tool-deaf response.
    effective_tier = tier
    if tier == ModelTier.TIER2 and not client._is_warm(ModelTier.TIER2):
        effective_tier = ModelTier.TIER3
        log.info("RunPod Tier2 cold — upgrading to Tier3 for full tool support")

    system_prompt = await assemble(
        user_id, content or "attachment", db=db, tier=effective_tier,
        user_lat=location_lat, user_lng=location_lng,
    )

    # Tools available for TIER2 and TIER3.
    # TIER2 (RunPod) ignores them in the payload but the Sonnet fallback path
    # uses them — so we must build them here regardless, otherwise a RunPod
    # timeout with a "warm" start would fall back to Sonnet with no tools.
    # TIER1 (Haiku) handles simple queries and never needs tool execution.
    tools = [
        CREATE_TASK_TOOL,
        CREATE_CALENDAR_EVENT_TOOL,
        PROPOSE_CALENDAR_EVENT_TOOL,
        PROPOSE_TASK_TOOL,
        SAVE_MEMORY_TOOL,
        SAVE_TO_SECOND_BRAIN_TOOL,
        READ_EMAIL_TOOL,
        SEND_EMAIL_TOOL,
        READ_MEETING_TOOL,
        SYNC_MEETINGS_TOOL,
        WEB_SEARCH_TOOL,
        GENERATE_DOCUMENT_TOOL,
        GENERATE_PRESENTATION_TOOL,
        GENERATE_PDF_TOOL,
        LOOKUP_CONTACT_TOOL,
        SEARCH_CONTACTS_TOOL,
        CREATE_CONTACT_TOOL,
        UPDATE_CONTACT_TOOL,
        SEARCH_PLACES_TOOL,
        SAVE_PLACE_TOOL,
        GET_SAVED_PLACES_TOOL,
        CREATE_AGENT_JOB_TOOL,
    ] if effective_tier != ModelTier.TIER1 else None
    queue: asyncio.Queue = asyncio.Queue()

    async def background_generate() -> None:
        full_response: list[str] = []
        # Card payloads emitted during streaming, persisted on the message so
        # they re-hydrate on conversation reload.
        tool_results: list[dict] = []
        model_used: str = ""
        tokens_used: int = 0

        try:
            # ── Phase 1: stream the model response ──────────────────────────
            async with AsyncSessionLocal() as bg_db:

                # Capture location for tool executor (search_places default near)
                _user_lat = location_lat
                _user_lng = location_lng

                async def _emit_card(event: dict) -> None:
                    """Emit a card event to the SSE queue AND record it for persistence."""
                    tool_results.append(event)
                    await queue.put(sse_event(event))

                async def _tool_executor(name: str, tool_input: dict) -> str:
                    if name == "create_task":
                        task = Task(
                            user_id=user_id,
                            title=tool_input["title"],
                            description=tool_input.get("description"),
                            priority=tool_input.get("priority", "normal"),
                            status="inbox",
                            source="chat",
                        )
                        bg_db.add(task)
                        await bg_db.commit()
                        priority = tool_input.get("priority", "normal")
                        return f"Task created: '{tool_input['title']}' added to inbox (priority: {priority})."

                    if name == "create_calendar_event":
                        try:
                            from sqlalchemy import select as _select
                            from db.models import Connector
                            conn_result = await bg_db.execute(
                                _select(Connector).where(
                                    Connector.user_id == user_id,
                                    Connector.name == "Google Calendar",
                                )
                            )
                            conn = conn_result.scalar_one_or_none()
                            if not conn or not conn.auth.get("refresh_token"):
                                return "Google Calendar not connected. Event not created."

                            from connectors.google_calendar import GoogleCalendarClient
                            from datetime import datetime, timedelta
                            import asyncio as _asyncio

                            gcal = GoogleCalendarClient(conn.auth)
                            start_dt = datetime.fromisoformat(tool_input["datetime_iso"])
                            duration = tool_input.get("duration_min", 60)
                            end_dt = start_dt + timedelta(minutes=duration)

                            event_body = {
                                "summary": tool_input["title"],
                                "start": {"dateTime": start_dt.isoformat(), "timeZone": str(start_dt.tzinfo or "UTC")},
                                "end": {"dateTime": end_dt.isoformat(), "timeZone": str(start_dt.tzinfo or "UTC")},
                            }
                            if tool_input.get("description"):
                                event_body["description"] = tool_input["description"]
                            if tool_input.get("location"):
                                event_body["location"] = tool_input["location"]
                            if tool_input.get("attendees"):
                                event_body["attendees"] = [{"email": e} for e in tool_input["attendees"]]

                            loop = _asyncio.get_event_loop()
                            result = await loop.run_in_executor(
                                None, lambda: gcal.create_event(**event_body)
                            )
                            link = result.get("htmlLink", "")
                            return f"Calendar event created: '{tool_input['title']}' on {start_dt.strftime('%b %-d at %-I:%M %p')}." + (f" View: {link}" if link else "")
                        except Exception as exc:
                            log.warning("create_calendar_event failed: %s", exc)
                            return f"Failed to create calendar event: {exc}"

                    if name == "save_memory":
                        try:
                            from memory import mnemon as _mnemon
                            await _mnemon.write(
                                bg_db,
                                user_id,
                                tool_input["content"],
                                domain=tool_input.get("domain", "work"),
                                source="chat",
                                importance=tool_input.get("importance", 3),
                            )
                            return f"Saved to memory: '{tool_input['content'][:80]}...'" if len(tool_input["content"]) > 80 else f"Saved to memory: '{tool_input['content']}'"
                        except Exception as exc:
                            log.warning("save_memory tool failed: %s", exc)
                            return f"Failed to save memory: {exc}"

                    if name == "save_to_second_brain":
                        try:
                            from memory import second_brain as _sb
                            await _sb.ingest_document(
                                bg_db,
                                user_id,
                                content=tool_input["content"],
                                title=tool_input["title"],
                                tags=tool_input.get("tags", []),
                                domain=tool_input.get("domain", "work"),
                            )
                            return f"Saved to Second Brain: '{tool_input['title']}'"
                        except Exception as exc:
                            log.warning("save_to_second_brain tool failed: %s", exc)
                            return f"Failed to save to Second Brain: {exc}"

                    if name in ("lookup_contact", "search_contacts"):
                        try:
                            from sqlalchemy import select as _select, or_ as _or, func as _func
                            from db.models import Contact, Connector
                            query = (tool_input.get("query") or "").strip()
                            limit  = int(tool_input.get("limit", 25)) if name == "search_contacts" else 5
                            offset = int(tool_input.get("offset", 0))

                            # Total unique contacts (by email) the user has
                            total_stmt = (
                                _select(_func.count(_func.distinct(Contact.primary_email)))
                                .where(Contact.user_id == user_id)
                            )
                            total_result = await bg_db.execute(total_stmt)
                            total_unique = total_result.scalar_one() or 0

                            if query:
                                needle = f"%{query}%"
                                filter_clause = _or(
                                    Contact.display_name.ilike(needle),
                                    Contact.primary_email.ilike(needle),
                                    Contact.organization.ilike(needle),
                                    Contact.primary_phone.ilike(needle),
                                )
                                stmt = (
                                    _select(Contact)
                                    .where(Contact.user_id == user_id, filter_clause)
                                    .order_by(Contact.is_other_contact, Contact.display_name)
                                    .offset(offset)
                                    .limit(limit * 4)  # over-fetch to allow email dedup
                                )
                            else:
                                # Browse mode — no filter, ordered by saved contacts first then name
                                stmt = (
                                    _select(Contact)
                                    .where(Contact.user_id == user_id)
                                    .order_by(Contact.is_other_contact, Contact.display_name)
                                    .offset(offset)
                                    .limit(limit * 4)
                                )

                            result = await bg_db.execute(stmt)
                            raw = result.scalars().all()

                            # Deduplicate by primary_email at query time
                            # (saved contact wins over other_contact for same email)
                            seen_emails: set[str] = set()
                            matches: list = []
                            for c in raw:
                                key = (c.primary_email or "").lower().strip() or c.id
                                if key in seen_emails:
                                    continue
                                seen_emails.add(key)
                                matches.append(c)
                                if len(matches) >= limit:
                                    break

                            # Fallback for lookup_contact: live Google search if no local hit
                            live_cards: list[dict] = []
                            if name == "lookup_contact" and not matches:
                                conn_result = await bg_db.execute(
                                    _select(Connector).where(
                                        Connector.user_id == user_id,
                                        Connector.name == "Google Contacts",
                                    )
                                )
                                conn = conn_result.scalar_one_or_none()
                                if conn and conn.auth.get("refresh_token"):
                                    from connectors.google_people import GooglePeopleClient
                                    import asyncio as _asyncio
                                    loop = _asyncio.get_event_loop()
                                    try:
                                        gp = GooglePeopleClient(conn.auth)
                                        live = await loop.run_in_executor(
                                            None, lambda: gp.search_contacts(query, page_size=5)
                                        )
                                        if live:
                                            lines = []
                                            for person in live:
                                                names  = person.get("names", [])
                                                emails = person.get("emailAddresses", [])
                                                phones = person.get("phoneNumbers", [])
                                                orgs   = person.get("organizations", [])
                                                name_  = (names[0].get("displayName") if names else None) or "Unknown"
                                                em     = emails[0].get("value") if emails else None
                                                ph     = phones[0].get("value") if phones else None
                                                org_   = orgs[0].get("name") if orgs else None
                                                title_ = orgs[0].get("title") if orgs else None
                                                parts  = [name_]
                                                if title_ or org_:
                                                    parts.append(f"({', '.join(p for p in [title_, org_] if p)})")
                                                if em:
                                                    parts.append(f"<{em}>")
                                                if ph:
                                                    parts.append(f"📞 {ph}")
                                                lines.append(" ".join(parts))
                                                live_cards.append({
                                                    "id": None,
                                                    "display_name": name_,
                                                    "primary_email": em,
                                                    "primary_phone": ph,
                                                    "organization": org_,
                                                    "job_title": title_,
                                                    "tars_context": None,
                                                    "source": "google_live",
                                                })
                                            if live_cards:
                                                await _emit_card({
                                                    "type": "contact_card",
                                                    "contacts": live_cards,
                                                })
                                            return f"Live Google search results for '{query}':\n" + "\n".join(lines)
                                    except Exception as exc:
                                        log.warning("Live People API search failed: %s", exc)

                            if not matches:
                                if query:
                                    return f"No contacts found matching '{query}'. Total contacts in database: {total_unique}."
                                return f"No contacts found. Total contacts in database: {total_unique}."

                            lines = []
                            cards: list[dict] = []
                            for c in matches:
                                parts = [c.display_name or "Unknown"]
                                if c.job_title or c.organization:
                                    parts.append(f"({', '.join(p for p in [c.job_title, c.organization] if p)})")
                                if c.primary_email:
                                    parts.append(f"<{c.primary_email}>")
                                if c.primary_phone:
                                    parts.append(f"📞 {c.primary_phone}")
                                line = " ".join(parts)
                                if c.tars_context:
                                    line += f"\n  context: {c.tars_context[:200]}"
                                lines.append(line)
                                cards.append({
                                    "id": c.id,
                                    "display_name": c.display_name,
                                    "primary_email": c.primary_email,
                                    "primary_phone": c.primary_phone,
                                    "organization": c.organization,
                                    "job_title": c.job_title,
                                    "tars_context": c.tars_context,
                                    "source": "local",
                                    "is_other_contact": c.is_other_contact,
                                })
                            if cards:
                                await _emit_card({
                                    "type": "contact_card",
                                    "contacts": cards,
                                })
                            if query:
                                header = f"Found {len(matches)} contact(s) matching '{query}' (total unique contacts: {total_unique}):"
                            else:
                                shown_range = f"{offset + 1}–{offset + len(matches)}"
                                header = f"Showing {shown_range} of {total_unique} unique contacts (use offset to page):"
                            return header + "\n" + "\n".join(lines)
                        except Exception as exc:
                            log.warning("%s tool failed: %s", name, exc)
                            return f"Contact lookup failed: {exc}"

                    if name == "create_contact":
                        try:
                            from sqlalchemy import select as _select
                            from db.models import Contact, Connector
                            from connectors.google_people import GooglePeopleClient, to_contact_dict
                            import asyncio as _asyncio

                            # Get Google Contacts connector
                            conn_result = await bg_db.execute(
                                _select(Connector).where(
                                    Connector.user_id == user_id,
                                    Connector.name == "Google Contacts",
                                )
                            )
                            conn = conn_result.scalar_one_or_none()
                            if not conn or not conn.auth.get("refresh_token"):
                                return "Google Contacts not connected — cannot create contact."

                            gp = GooglePeopleClient(conn.auth)
                            loop = _asyncio.get_event_loop()
                            person = await loop.run_in_executor(None, lambda: gp.create_contact(
                                name=tool_input["name"],
                                email=tool_input.get("email"),
                                phone=tool_input.get("phone"),
                                organization=tool_input.get("organization"),
                                job_title=tool_input.get("job_title"),
                                biography=tool_input.get("notes"),
                            ))

                            # Sync the new contact into the local DB immediately
                            contact_data = to_contact_dict(person)
                            contact_data["is_other_contact"] = False
                            from datetime import datetime as _dt, timezone as _tz
                            contact_data["last_synced_at"] = _dt.now(_tz.utc)
                            resource_name = contact_data.get("google_resource_name")
                            existing = None
                            if resource_name:
                                ex_result = await bg_db.execute(
                                    _select(Contact).where(Contact.google_resource_name == resource_name)
                                )
                                existing = ex_result.scalar_one_or_none()
                            if existing:
                                for k, v in contact_data.items():
                                    setattr(existing, k, v)
                            else:
                                bg_db.add(Contact(user_id=user_id, **contact_data))
                            await bg_db.commit()

                            parts = [f"Contact created: {tool_input['name']}"]
                            if tool_input.get("email"):
                                parts.append(f"<{tool_input['email']}>")
                            if tool_input.get("phone"):
                                parts.append(f"📞 {tool_input['phone']}")
                            if tool_input.get("organization"):
                                parts.append(f"@ {tool_input['organization']}")
                            return " ".join(parts) + ". Saved to Google Contacts and local database."
                        except Exception as exc:
                            log.warning("create_contact tool failed: %s", exc)
                            return f"Failed to create contact: {exc}"

                    if name == "update_contact":
                        try:
                            from sqlalchemy import select as _select, or_ as _or
                            from db.models import Contact, Connector
                            from connectors.google_people import GooglePeopleClient, to_contact_dict
                            import asyncio as _asyncio

                            query = (tool_input.get("query") or "").strip()
                            if not query:
                                return "Provide a name or email to identify the contact to update."

                            # Find the contact in local DB (saved contacts only — need resource_name + etag)
                            needle = f"%{query}%"
                            stmt = (
                                _select(Contact)
                                .where(
                                    Contact.user_id == user_id,
                                    Contact.is_other_contact == False,  # noqa: E712
                                    _or(
                                        Contact.display_name.ilike(needle),
                                        Contact.primary_email.ilike(needle),
                                    ),
                                )
                                .limit(1)
                            )
                            result = await bg_db.execute(stmt)
                            contact = result.scalar_one_or_none()

                            if not contact or not contact.google_resource_name:
                                return (
                                    f"No saved contact found matching '{query}'. "
                                    "If this is an unsaved 'other contact', create them first with create_contact."
                                )

                            # Determine which Google fields to update
                            update_fields: dict = {}
                            person_fields_list: list[str] = []

                            if tool_input.get("name"):
                                update_fields["names"] = [{"unstructuredName": tool_input["name"]}]
                                person_fields_list.append("names")
                            if tool_input.get("email"):
                                update_fields["emailAddresses"] = [{"value": tool_input["email"]}]
                                person_fields_list.append("emailAddresses")
                            if tool_input.get("phone"):
                                update_fields["phoneNumbers"] = [{"value": tool_input["phone"]}]
                                person_fields_list.append("phoneNumbers")
                            if tool_input.get("organization") or tool_input.get("job_title"):
                                org_entry: dict = {}
                                if tool_input.get("organization"):
                                    org_entry["name"] = tool_input["organization"]
                                elif contact.organization:
                                    org_entry["name"] = contact.organization
                                if tool_input.get("job_title"):
                                    org_entry["title"] = tool_input["job_title"]
                                elif contact.job_title:
                                    org_entry["title"] = contact.job_title
                                update_fields["organizations"] = [org_entry]
                                person_fields_list.append("organizations")
                            if tool_input.get("notes"):
                                update_fields["biographies"] = [{"value": tool_input["notes"], "contentType": "TEXT_PLAIN"}]
                                person_fields_list.append("biographies")

                            if not update_fields:
                                return "No fields provided to update. Specify at least one of: name, email, phone, organization, job_title, notes."

                            conn_result = await bg_db.execute(
                                _select(Connector).where(
                                    Connector.user_id == user_id,
                                    Connector.name == "Google Contacts",
                                )
                            )
                            conn = conn_result.scalar_one_or_none()
                            if not conn or not conn.auth.get("refresh_token"):
                                return "Google Contacts not connected — cannot update contact."

                            gp = GooglePeopleClient(conn.auth)
                            loop = _asyncio.get_event_loop()
                            updated_person = await loop.run_in_executor(None, lambda: gp.update_contact(
                                resource_name=contact.google_resource_name,
                                etag=contact.etag or "*",
                                fields=update_fields,
                                update_person_fields=",".join(person_fields_list),
                            ))

                            # Update local DB row from the returned person
                            updated_data = to_contact_dict(updated_person)
                            from datetime import datetime as _dt, timezone as _tz
                            updated_data["last_synced_at"] = _dt.now(_tz.utc)
                            for k, v in updated_data.items():
                                if hasattr(contact, k):
                                    setattr(contact, k, v)
                            await bg_db.commit()

                            changed = []
                            if tool_input.get("name"):     changed.append(f"name → {tool_input['name']}")
                            if tool_input.get("phone"):    changed.append(f"phone → {tool_input['phone']}")
                            if tool_input.get("email"):    changed.append(f"email → {tool_input['email']}")
                            if tool_input.get("organization"): changed.append(f"org → {tool_input['organization']}")
                            if tool_input.get("job_title"): changed.append(f"title → {tool_input['job_title']}")
                            if tool_input.get("notes"):    changed.append("notes updated")
                            display = contact.display_name or query
                            return f"Updated {display}: {', '.join(changed)}. Synced to Google Contacts."
                        except Exception as exc:
                            log.warning("update_contact tool failed: %s", exc)
                            return f"Failed to update contact: {exc}"

                    if name == "read_email":
                        try:
                            from sqlalchemy import select as _select
                            from db.models import Connector
                            conn_result = await bg_db.execute(
                                _select(Connector).where(
                                    Connector.user_id == user_id,
                                    Connector.name == "Gmail",
                                )
                            )
                            conn = conn_result.scalar_one_or_none()
                            if not conn or not conn.auth.get("refresh_token"):
                                return "Gmail not connected."

                            from connectors.gmail import GmailClient, extract_thread_text
                            import asyncio as _asyncio

                            gclient = GmailClient(conn.auth)
                            loop = _asyncio.get_event_loop()

                            thread_id = tool_input.get("thread_id", "").strip()
                            if not thread_id and tool_input.get("search_query"):
                                threads = await loop.run_in_executor(
                                    None,
                                    lambda: gclient.list_threads(
                                        query=tool_input["search_query"], max_results=1
                                    ),
                                )
                                if not threads:
                                    return "No email found matching that search."
                                thread_id = threads[0]["id"]

                            if not thread_id:
                                return "Provide a thread_id or search_query."

                            # thread_id from context is first 8 chars — find full id
                            if len(thread_id) == 8:
                                candidates = await loop.run_in_executor(
                                    None,
                                    lambda: gclient.list_threads(query="in:inbox", max_results=20),
                                )
                                match = next((t["id"] for t in candidates if t["id"].startswith(thread_id)), None)
                                if match:
                                    thread_id = match

                            thread = await loop.run_in_executor(None, lambda: gclient.get_thread(thread_id))
                            body = extract_thread_text(thread)
                            if not body:
                                return "Email body is empty or could not be extracted."

                            # Enqueue any unique sender as a PendingContact for review
                            try:
                                senders: list[str] = []
                                subject = ""
                                for msg in thread.get("messages", []):
                                    headers = {h["name"]: h["value"] for h in msg.get("payload", {}).get("headers", [])}
                                    if not subject:
                                        subject = headers.get("Subject", "")
                                    s = headers.get("From", "").strip()
                                    if s and s not in senders:
                                        senders.append(s)
                                if senders:
                                    from jobs.pending_contacts import enqueue_from_strings
                                    context = f"Emailed about: {subject[:120]}" if subject else "Detected from email thread"
                                    new_count = await enqueue_from_strings(
                                        bg_db, user_id, senders,
                                        source="email",
                                        source_id=thread_id,
                                        extracted_context=context,
                                    )
                                    if new_count:
                                        await bg_db.commit()
                                        log.info("read_email: queued %d new contact(s) for review", new_count)
                            except Exception as exc:
                                log.warning("Failed to enqueue email senders: %s", exc)

                            # Cap at 4000 chars to stay within token budget
                            if len(body) > 4000:
                                body = body[:4000] + "\n\n[... email truncated ...]"
                            return body
                        except Exception as exc:
                            log.warning("read_email tool failed: %s", exc)
                            return f"Failed to read email: {exc}"

                    if name == "send_email":
                        try:
                            from sqlalchemy import select as _select
                            from db.models import Connector
                            conn_result = await bg_db.execute(
                                _select(Connector).where(
                                    Connector.user_id == user_id,
                                    Connector.name == "Gmail",
                                )
                            )
                            conn = conn_result.scalar_one_or_none()
                            if not conn or not conn.auth.get("refresh_token"):
                                return "Gmail not connected — cannot send email."

                            from connectors.gmail import GmailClient
                            import asyncio as _asyncio

                            gclient = GmailClient(conn.auth)
                            loop = _asyncio.get_event_loop()
                            result = await loop.run_in_executor(
                                None,
                                lambda: gclient.send_email(
                                    to=tool_input["to"],
                                    subject=tool_input["subject"],
                                    body=tool_input["body"],
                                    cc=tool_input.get("cc"),
                                    thread_id=tool_input.get("thread_id"),
                                ),
                            )
                            return result
                        except Exception as exc:
                            log.warning("send_email tool failed: %s", exc)
                            return f"Failed to send email: {exc}"

                    if name == "sync_meetings":
                        try:
                            from core.config import settings as _settings
                            if not _settings.fireflies_api_key:
                                return "Fireflies API key not configured."
                            from connectors.fireflies import FirefliesClient
                            from jobs.meeting_processor import ingest_from_webhook, process_meeting as _proc_meeting
                            ff_client = FirefliesClient(_settings.fireflies_api_key)
                            transcripts = await ff_client.list_recent(limit=20)
                            synced = 0
                            skipped = 0
                            for t in transcripts:
                                tid = t.get("id")
                                if not tid:
                                    continue
                                new_id = await ingest_from_webhook(bg_db, user_id, tid)
                                if new_id:
                                    await _proc_meeting(bg_db, new_id, user_id)
                                    synced += 1
                                else:
                                    skipped += 1
                            if synced:
                                return f"Synced {synced} new meeting{'s' if synced != 1 else ''} from Fireflies ({skipped} already up to date)."
                            else:
                                return f"All {skipped} recent Fireflies meetings are already up to date."
                        except Exception as exc:
                            log.warning("sync_meetings tool failed: %s", exc)
                            return f"Failed to sync meetings: {exc}"

                    if name == "read_meeting":
                        try:
                            from sqlalchemy import select as _select
                            from db.models import Meeting, MeetingActionItem

                            meeting_id = tool_input.get("meeting_id", "").strip()
                            include_transcript = tool_input.get("include_transcript", False)

                            m_result = await bg_db.execute(
                                _select(Meeting).where(
                                    Meeting.id == meeting_id,
                                    Meeting.user_id == user_id,
                                )
                            )
                            meeting = m_result.scalar_one_or_none()
                            if not meeting:
                                return f"Meeting '{meeting_id}' not found."

                            ai_result = await bg_db.execute(
                                _select(MeetingActionItem).where(
                                    MeetingActionItem.meeting_id == meeting_id
                                )
                            )
                            action_items = ai_result.scalars().all()

                            parts = [f"# {meeting.title}"]
                            if meeting.started_at:
                                parts.append(f"Date: {meeting.started_at.strftime('%B %-d, %Y')}")
                            if meeting.attendees:
                                parts.append(f"Attendees: {', '.join(meeting.attendees)}")
                            parts.append(f"Status: {meeting.status}")

                            if meeting.summary:
                                parts.append(f"\n## Summary\n{meeting.summary}")

                            if action_items:
                                parts.append(f"\n## Action Items ({len(action_items)})")
                                for ai in action_items:
                                    owner = f" [{ai.owner}]" if ai.owner else ""
                                    done = " ✓" if ai.task_id else ""
                                    parts.append(f"  • {ai.raw_text}{owner}{done}")

                            if include_transcript and meeting.transcript:
                                # Cap at 20k chars — enough for a 2h meeting without token explosion
                                transcript = meeting.transcript
                                if len(transcript) > 20000:
                                    transcript = transcript[:20000] + "\n\n[... transcript truncated — full text available in Meetings section ...]"
                                parts.append(f"\n## Transcript\n{transcript}")
                            elif not include_transcript and meeting.transcript:
                                parts.append(f"\n(Transcript available — call read_meeting again with include_transcript=true to see it.)")

                            return "\n".join(parts)
                        except Exception as exc:
                            log.warning("read_meeting tool failed: %s", exc)
                            return f"Failed to read meeting: {exc}"

                    if name == "web_search":
                        try:
                            from core.config import settings as _settings
                            if not _settings.tavily_api_key:
                                return "Web search is not configured (missing TAVILY_API_KEY)."
                            import httpx as _httpx
                            query = tool_input["query"]
                            depth = tool_input.get("search_depth", "basic")
                            async with _httpx.AsyncClient(timeout=15.0) as _hx:
                                resp = await _hx.post(
                                    "https://api.tavily.com/search",
                                    json={
                                        "api_key": _settings.tavily_api_key,
                                        "query": query,
                                        "search_depth": depth,
                                        "max_results": 6,
                                        "include_answer": True,
                                    },
                                )
                                resp.raise_for_status()
                                data = resp.json()

                            lines = [f"Search: {query}\n"]
                            if data.get("answer"):
                                lines.append(f"Summary: {data['answer']}\n")
                            for r in data.get("results", []):
                                lines.append(f"• {r['title']}")
                                lines.append(f"  {r['url']}")
                                if r.get("content"):
                                    lines.append(f"  {r['content'][:300]}")
                                lines.append("")
                            return "\n".join(lines)
                        except Exception as exc:
                            log.warning("web_search tool failed: %s", exc)
                            return f"Web search failed: {exc}"

                    if name == "generate_document":
                        try:
                            import re as _re
                            from io import BytesIO
                            import docx as _docx
                            title = tool_input.get("title", "Document")
                            content_md = tool_input.get("content", "")
                            fn_base = tool_input.get("filename") or _re.sub(r"[^\w\s\-]", "", title).strip().replace(" ", "_")[:50]
                            doc = _docx.Document()
                            doc.add_heading(title, 0)
                            for line in content_md.split("\n"):
                                s = line.strip()
                                if not s:
                                    continue
                                if s.startswith("### "):
                                    doc.add_heading(s[4:], level=3)
                                elif s.startswith("## "):
                                    doc.add_heading(s[3:], level=2)
                                elif s.startswith("# "):
                                    doc.add_heading(s[2:], level=1)
                                elif s.startswith("- ") or s.startswith("* "):
                                    doc.add_paragraph(s[2:], style="List Bullet")
                                elif _re.match(r"^\d+\.", s):
                                    doc.add_paragraph(_re.sub(r"^\d+\.\s*", "", s), style="List Number")
                                else:
                                    doc.add_paragraph(s)
                            buf = BytesIO()
                            doc.save(buf)
                            raw = buf.getvalue()
                            b64 = base64.b64encode(raw).decode()
                            filename = fn_base + ".docx"
                            artifact = Artifact(
                                user_id=user_id, filename=filename, type="document",
                                source="chat", source_id=conversation_id,
                                content="base64:" + b64, version=1,
                                size_bytes=len(raw), tags=["generated"],
                            )
                            bg_db.add(artifact)
                            await bg_db.commit()
                            await bg_db.refresh(artifact)
                            await _emit_card({"type": "artifact_created", "artifact_id": artifact.id, "filename": filename, "filetype": "docx"})
                            return f"Generated '{filename}' ({len(raw):,} bytes). Saved to Artifacts."
                        except Exception as exc:
                            log.warning("generate_document failed: %s", exc)
                            return f"Failed to generate document: {exc}"

                    if name == "generate_presentation":
                        try:
                            import re as _re
                            from io import BytesIO
                            from pptx import Presentation as _Prs
                            from pptx.util import Pt as _Pt
                            title = tool_input.get("title", "Presentation")
                            subtitle = tool_input.get("subtitle", "")
                            slides_data = tool_input.get("slides", [])
                            fn_base = tool_input.get("filename") or _re.sub(r"[^\w\s\-]", "", title).strip().replace(" ", "_")[:50]
                            prs = _Prs()
                            # Title slide
                            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
                            title_slide.shapes.title.text = title
                            if subtitle and len(title_slide.placeholders) > 1:
                                title_slide.placeholders[1].text = subtitle
                            # Content slides
                            for s in slides_data:
                                slide = prs.slides.add_slide(prs.slide_layouts[1])
                                slide.shapes.title.text = s.get("title", "")
                                tf = slide.placeholders[1].text_frame
                                tf.clear()
                                bullets = s.get("bullets", [])
                                for i, bullet in enumerate(bullets):
                                    if i == 0:
                                        tf.text = bullet
                                    else:
                                        p = tf.add_paragraph()
                                        p.text = bullet
                                        p.level = 0
                            buf = BytesIO()
                            prs.save(buf)
                            raw = buf.getvalue()
                            b64 = base64.b64encode(raw).decode()
                            filename = fn_base + ".pptx"
                            artifact = Artifact(
                                user_id=user_id, filename=filename, type="document",
                                source="chat", source_id=conversation_id,
                                content="base64:" + b64, version=1,
                                size_bytes=len(raw), tags=["generated", "presentation"],
                            )
                            bg_db.add(artifact)
                            await bg_db.commit()
                            await bg_db.refresh(artifact)
                            await _emit_card({"type": "artifact_created", "artifact_id": artifact.id, "filename": filename, "filetype": "pptx"})
                            return f"Generated '{filename}' with {len(slides_data) + 1} slides ({len(raw):,} bytes). Saved to Artifacts."
                        except Exception as exc:
                            log.warning("generate_presentation failed: %s", exc)
                            return f"Failed to generate presentation: {exc}"

                    if name == "generate_pdf":
                        try:
                            import re as _re
                            from io import BytesIO
                            from reportlab.lib.pagesizes import A4
                            from reportlab.lib.styles import getSampleStyleSheet
                            from reportlab.lib.units import inch
                            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
                            title = tool_input.get("title", "Document")
                            content_md = tool_input.get("content", "")
                            fn_base = tool_input.get("filename") or _re.sub(r"[^\w\s\-]", "", title).strip().replace(" ", "_")[:50]
                            buf = BytesIO()
                            doc = SimpleDocTemplate(buf, pagesize=A4,
                                rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=72)
                            styles = getSampleStyleSheet()
                            story = [Paragraph(title, styles["Title"]), Spacer(1, 0.2 * inch)]
                            for line in content_md.split("\n"):
                                s = line.strip()
                                if not s:
                                    story.append(Spacer(1, 0.08 * inch))
                                    continue
                                if s.startswith("### "):
                                    story.append(Paragraph(s[4:], styles["Heading3"]))
                                elif s.startswith("## "):
                                    story.append(Paragraph(s[3:], styles["Heading2"]))
                                elif s.startswith("# "):
                                    story.append(Paragraph(s[2:], styles["Heading1"]))
                                elif s.startswith("- ") or s.startswith("* "):
                                    story.append(Paragraph("&#8226; " + s[2:], styles["Normal"]))
                                else:
                                    story.append(Paragraph(s, styles["Normal"]))
                                    story.append(Spacer(1, 0.05 * inch))
                            doc.build(story)
                            raw = buf.getvalue()
                            b64 = base64.b64encode(raw).decode()
                            filename = fn_base + ".pdf"
                            artifact = Artifact(
                                user_id=user_id, filename=filename, type="document",
                                source="chat", source_id=conversation_id,
                                content="base64:" + b64, version=1,
                                size_bytes=len(raw), tags=["generated", "pdf"],
                            )
                            bg_db.add(artifact)
                            await bg_db.commit()
                            await bg_db.refresh(artifact)
                            await _emit_card({"type": "artifact_created", "artifact_id": artifact.id, "filename": filename, "filetype": "pdf"})
                            return f"Generated '{filename}' ({len(raw):,} bytes). Saved to Artifacts."
                        except Exception as exc:
                            log.warning("generate_pdf failed: %s", exc)
                            return f"Failed to generate PDF: {exc}"

                    # ── Places tools ──────────────────────────────────────────────────
                    if name == "search_places":
                        try:
                            from connectors.places import PlacesClient
                            import asyncio as _asyncio

                            query    = (tool_input.get("query") or "").strip()
                            near     = tool_input.get("near")
                            category = tool_input.get("category")
                            limit    = int(tool_input.get("limit") or 5)

                            # Fall back to user's current location when no "near" was specified
                            user_coords_str = (
                                f"{_user_lat},{_user_lng}"
                                if _user_lat is not None and _user_lng is not None
                                else None
                            )
                            effective_near = near or user_coords_str

                            client_p = PlacesClient()
                            loop = _asyncio.get_event_loop()

                            # ── "Where am I?" / reverse-geocode intent ──────────
                            _location_queries = {
                                "my location", "current location", "where am i", "where are we",
                                "my position", "here", "current position",
                            }
                            _is_location_query = query.lower().strip("?").strip() in _location_queries

                            if _is_location_query and _user_lat is not None and _user_lng is not None:
                                geo = await loop.run_in_executor(
                                    None, lambda: client_p.reverse_geocode(_user_lat, _user_lng)
                                )
                                if geo:
                                    results = [geo]
                                else:
                                    results = []
                            elif category and _user_lat is not None and _user_lng is not None and not near:
                                # We have exact coords — go straight to Overpass nearby search
                                results = await loop.run_in_executor(
                                    None, lambda: client_p.search_nearby(_user_lat, _user_lng, category, limit=limit)
                                )
                            elif category and effective_near and effective_near != user_coords_str:
                                # Named location — geocode it first, then nearby
                                geo = await loop.run_in_executor(None, lambda: client_p.search(effective_near, limit=1))
                                if geo:
                                    results = await loop.run_in_executor(
                                        None, lambda: client_p.search_nearby(geo[0]["lat"], geo[0]["lng"], category, limit=limit)
                                    )
                                else:
                                    results = await loop.run_in_executor(
                                        None, lambda: client_p.search(f"{category} {effective_near}", limit=limit)
                                    )
                            elif not category and _user_lat is not None and _user_lng is not None and not near:
                                # No category, no named location, but have GPS — do nearby text search biased by coords
                                results = await loop.run_in_executor(
                                    None, lambda: client_p.search(query, near=user_coords_str, limit=limit)
                                )
                            else:
                                results = await loop.run_in_executor(
                                    None, lambda: client_p.search(query, near=effective_near, limit=limit)
                                )

                            if not results:
                                return f"No places found for '{query}'{' near ' + near if near else ''}."

                            # Emit place_card SSE event
                            cards = [
                                {
                                    "name":         p.get("name", ""),
                                    "address":      p.get("address") or p.get("display_name", ""),
                                    "lat":          p.get("lat"),
                                    "lng":          p.get("lng"),
                                    "category":     p.get("category"),
                                    "osm_id":       p.get("osm_id"),
                                    "osm_type":     p.get("osm_type"),
                                    "source":       "osm",
                                    "is_saved":     False,
                                }
                                for p in results
                            ]
                            await _emit_card({"type": "place_card", "places": cards})

                            lines = []
                            for p in results:
                                addr = p.get("address") or p.get("display_name", "")
                                lines.append(f"• {p['name']}" + (f" — {addr}" if addr else ""))
                            return f"Found {len(results)} place(s):\n" + "\n".join(lines)

                        except Exception as exc:
                            log.warning("search_places tool failed: %s", exc)
                            return f"Place search failed: {exc}"

                    if name == "save_place":
                        try:
                            from db.models import Place as _Place

                            p = _Place(
                                user_id=user_id,
                                name=tool_input["name"],
                                address=tool_input.get("address"),
                                lat=float(tool_input["lat"]),
                                lng=float(tool_input["lng"]),
                                category=tool_input.get("category"),
                                tags=tool_input.get("tags") or [],
                                notes=tool_input.get("notes"),
                                osm_id=tool_input.get("osm_id"),
                                osm_type=tool_input.get("osm_type"),
                                source="manual" if not tool_input.get("osm_id") else "osm",
                                is_saved=True,
                            )
                            bg_db.add(p)
                            await bg_db.commit()

                            # Emit card for the saved place
                            await _emit_card({
                                "type": "place_card",
                                "places": [{
                                    "name":     p.name,
                                    "address":  p.address,
                                    "lat":      p.lat,
                                    "lng":      p.lng,
                                    "category": p.category,
                                    "osm_id":   p.osm_id,
                                    "osm_type": p.osm_type,
                                    "source":   p.source,
                                    "is_saved": True,
                                    "notes":    p.notes,
                                    "tags":     p.tags,
                                }],
                            })

                            msg = f"Saved '{p.name}'"
                            if p.address:
                                msg += f" ({p.address})"
                            if p.notes:
                                msg += f" — {p.notes}"
                            return msg + " to your places."

                        except Exception as exc:
                            log.warning("save_place tool failed: %s", exc)
                            return f"Failed to save place: {exc}"

                    if name == "get_saved_places":
                        try:
                            from sqlalchemy import select as _select, or_ as _or
                            from db.models import Place as _Place

                            query    = (tool_input.get("query") or "").strip()
                            category = (tool_input.get("category") or "").strip()
                            limit    = int(tool_input.get("limit") or 20)

                            stmt = _select(_Place).where(
                                _Place.user_id == user_id,
                                _Place.is_saved == True,  # noqa: E712
                            )
                            if query:
                                needle = f"%{query}%"
                                stmt = stmt.where(_or(
                                    _Place.name.ilike(needle),
                                    _Place.address.ilike(needle),
                                    _Place.notes.ilike(needle),
                                ))
                            if category:
                                stmt = stmt.where(_Place.category.ilike(f"%{category}%"))

                            stmt = stmt.order_by(_Place.created_at.desc()).limit(limit)
                            result = await bg_db.execute(stmt)
                            places = result.scalars().all()

                            if not places:
                                msg = "No saved places"
                                if query:
                                    msg += f" matching '{query}'"
                                if category:
                                    msg += f" in category '{category}'"
                                return msg + "."

                            cards = [
                                {
                                    "name":     pl.name,
                                    "address":  pl.address,
                                    "lat":      pl.lat,
                                    "lng":      pl.lng,
                                    "category": pl.category,
                                    "osm_id":   pl.osm_id,
                                    "osm_type": pl.osm_type,
                                    "source":   pl.source,
                                    "is_saved": True,
                                    "notes":    pl.notes,
                                    "tags":     pl.tags,
                                }
                                for pl in places
                            ]
                            await _emit_card({"type": "place_card", "places": cards})

                            lines = [f"• {pl.name}" + (f" — {pl.address}" if pl.address else "") for pl in places]
                            return f"Your saved places ({len(places)}):\n" + "\n".join(lines)

                        except Exception as exc:
                            log.warning("get_saved_places tool failed: %s", exc)
                            return f"Failed to retrieve saved places: {exc}"

                    if name == "create_agent_job":
                        try:
                            from db.models import AgentJob
                            agent_type = tool_input.get("agent_type", "evolutionarist")
                            instruction = tool_input.get("instruction", "")
                            job = AgentJob(
                                user_id=user_id,
                                agent_type=agent_type,
                                type="agent",
                                instruction=instruction,
                                repo_path="/Users/mike/Documents/TARS",
                                branch="dev",
                                status="pending",
                                conversation_id=conversation_id,
                            )
                            bg_db.add(job)
                            await bg_db.commit()
                            await bg_db.refresh(job)
                            # Start the agent in background
                            from agents.job_manager import start_job as _start_job
                            from db.session import AsyncSessionLocal as _ASL
                            import asyncio as _asyncio
                            _asyncio.create_task(_start_job(job.id, _ASL))
                            # Emit card event so frontend can show a link to the job
                            await _emit_card({
                                "type": "agent_job_created",
                                "job_id": job.id,
                                "agent_type": agent_type,
                                "instruction": instruction,
                            })
                            return f"Agent job created (ID: {job.id}). Type: {agent_type}. The agent is now running — watch the live stream in Agent Jobs."
                        except Exception as exc:
                            log.warning("create_agent_job tool failed: %s", exc)
                            return f"Failed to create agent job: {exc}"

                    return "Action completed."

                async for event in client.stream(messages, effective_tier, system=system_prompt, tools=tools, tool_executor=_tool_executor):
                    if event["type"] == "chunk":
                        full_response.append(event.get("text", ""))
                        await queue.put(sse_event(event))
                    elif event["type"] in ("calendar_suggest", "task_suggest", "contact_card", "place_card", "artifact_created"):
                        tool_results.append(event)
                        await queue.put(sse_event(event))
                    elif event["type"] == "done":
                        model_used = event.get("model", "")
                        tokens_used = event.get("tokens", 0)
                    elif event["type"] == "error":
                        await queue.put(sse_event(event))

            assistant_content = "".join(full_response)

            # If the model only emitted tool calls (no text) but did produce
            # cards, give the message a tiny placeholder so the bubble isn't
            # blank on reload. Cards re-render from tool_results.
            if not assistant_content.strip() and tool_results:
                kinds = {r.get("type") for r in tool_results}
                if "place_card" in kinds:
                    assistant_content = "Here's what I found:"
                elif "contact_card" in kinds:
                    assistant_content = "Here's who I found:"
                elif "artifact_created" in kinds:
                    assistant_content = "Done — see the file above."
                else:
                    assistant_content = "Done."

            # ── Phase 2: signal completion to client immediately ─────────────
            # Do NOT wait for DB saves / title gen / memory — client gets the
            # response right away, then we persist in the background.
            await queue.put(sse_event({"type": "done", "model": model_used, "tokens": tokens_used}))
            await queue.put(sse_done())

            # ── Phase 3: persist (client already showing the response) ───────
            async with AsyncSessionLocal() as bg_db:
                assistant_msg = Message(
                    conversation_id=conversation_id,
                    role="assistant",
                    content=assistant_content,
                    model_used=model_used or tier.value,
                    tokens_used=tokens_used,
                    tool_results=tool_results,
                )
                bg_db.add(assistant_msg)
                await bg_db.commit()

                title_msgs = messages + [{"role": "assistant", "content": assistant_content}]
                new_title = await _generate_title(title_msgs, client)
                if new_title:
                    await bg_db.execute(
                        sa_update(Conversation)
                        .where(Conversation.id == conversation_id)
                        .values(title=new_title)
                    )
                    await bg_db.commit()

                user_input = content or "[attachment]"
                await _extract_and_save_facts(bg_db, user_id, user_input, assistant_content, client)

        except Exception as exc:
            log.error("background_generate failed: %s", exc)
            await queue.put(sse_event({"type": "error", "error": str(exc)}))
        finally:
            await queue.put(None)

    task = asyncio.create_task(background_generate())
    _active_bg_tasks.add(task)
    task.add_done_callback(_active_bg_tasks.discard)

    async def generate():
        while True:
            item = await queue.get()
            if item is None:
                break
            yield item

    return StreamingResponse(
        generate(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        },
    )
