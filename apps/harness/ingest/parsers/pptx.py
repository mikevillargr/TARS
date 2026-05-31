"""PPTX parser using python-pptx."""
from __future__ import annotations
import io


def extract(content_bytes: bytes) -> str:
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        return "[PPTX parser unavailable — install python-pptx]"

    prs = Presentation(io.BytesIO(content_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        lines = [f"## Slide {i}"]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    lines.append(text)
        if len(lines) > 1:
            slides.append("\n".join(lines))
    return "\n\n".join(slides)
